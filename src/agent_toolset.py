

import json
import os
import logging
import requests
import random
from pathlib import Path
from datetime import datetime
from typing import Dict, Any

from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from openpyxl import Workbook
from openpyxl.chart import PieChart, BarChart, Reference
from openpyxl.styles import PatternFill, Font

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

class DocumentCreatorToolset:
    def __init__(self):
        self.openai_key = os.getenv("OPENAI_API_KEY")
        self.pexels_key = os.getenv("PEXELS_API_KEY")
        self.client = AsyncOpenAI(api_key=self.openai_key)
        self.model_name = "gpt-4o"
        
        self.supported_formats = ["pptx_slide_deck", "xlsx_spreadsheet"]
        self.last_brief = None
        self.last_plan = None
        
        self.image_cache_dir = "image_cache"
        os.makedirs(self.image_cache_dir, exist_ok=True)

    async def generate_document(self, brief: str) -> str:
        logger.info(f"📄 GENERATE_DOCUMENT CALLED: {brief}")
        try:
            self.last_brief = brief
            plan = await self._detect_format_and_plan(brief)
            fmt = plan.get("detected_format", "pptx_slide_deck")
            
            if fmt == "xlsx_spreadsheet":
                logger.info("📊 Generating Spreadsheet...")
                plan = await self._generate_full_xlsx_plan_from_brief(brief)
                doc_path = self._render_xlsx_from_plan(plan)
            else:
                logger.info("🖼️ Generating Presentation...")
                plan = await self._generate_full_ppt_plan_from_brief(brief)
                doc_path = self._render_ppt_from_plan(plan)

            self.last_plan = plan
            return json.dumps({"status": "success", "file_path": doc_path, "message": f"Successfully generated {fmt}"})
        except Exception as e:
            logger.exception("Error in generate_document")
            return json.dumps({"status": "error", "message": str(e)})

    async def _detect_format_and_plan(self, brief: str) -> Dict[str, Any]:
        prompt = f'Analyze "{brief}". Return JSON: {{ "detected_format": "pptx_slide_deck" }} or "xlsx_spreadsheet".'
        response = await self.client.chat.completions.create(
            model=self.model_name, messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}, temperature=0.1
        )
        return json.loads(response.choices[0].message.content)

    async def _generate_full_ppt_plan_from_brief(self, brief: str) -> Dict[str, Any]:
        prompt = f'''
        Create a detailed 10-slide presentation plan for: "{brief}".
        
        STRICT CONTENT RULES:
        1. SECTION 0 (TITLE SLIDE): Generate a professional, catchy, and creative title (DO NOT repeat the prompt) and a sophisticated subtitle.
        2. CONTENT SECTIONS (1-9): For EVERY slide, you MUST provide 5 to 6 DETAILED bullet points.
        3. TEXT QUALITY: Each bullet point must be a full, informative sentence (15-20 words). DO NOT use short phrases or simple headings.
        4. IMAGE QUERIES: Provide an aesthetic Pexels query for every slide.
        
        Return ONLY a JSON object:
        {{
          "detected_format": "pptx_slide_deck",
          "content_plan": {{
            "sections": [
              {{ "title": "Creative Title", "description": "Subtitle", "image_search_query": "term", "elements": [] }},
              {{ "title": "Slide Title", "description": "Context", "image_search_query": "term", "elements": [ {{"description": "Detailed Sentence 1..."}}, ... ] }}
            ]
          }}
        }}
        '''
        response = await self.client.chat.completions.create(
            model=self.model_name, messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}, temperature=0.7, max_tokens=4000
        )
        return json.loads(response.choices[0].message.content)

    async def _generate_full_xlsx_plan_from_brief(self, brief: str) -> Dict[str, Any]:
        prompt = f'Design a real-data spreadsheet for: {brief}. 4 sheets. RAW numbers. JSON ONLY.'
        response = await self.client.chat.completions.create(
            model=self.model_name, messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}, temperature=0.2 
        )
        return json.loads(response.choices[0].message.content)

    def _fetch_image_for_query(self, query: str) -> str:
        # CLEAN SYNTAX: Fixed the parenthesis/bracket closing from previous error
        safe_name = "".join([c if c.isalnum() or c in "-_" else "_" for c in query.lower()[:30]])
        cached_path = Path(self.image_cache_dir) / f"{safe_name}_{random.randint(1,100)}.jpg"
        headers = {"Authorization": self.pexels_key}
        params = {"query": query, "per_page": 5} 
        try:
            resp = requests.get("https://api.pexels.com/v1/search", headers=headers, params=params, timeout=10)
            if resp.status_code == 200 and resp.json().get("photos"):
                chosen_photo = resp.json()["photos"][0]
                img_resp = requests.get(chosen_photo["src"]["large"], timeout=10)
                with open(cached_path, "wb") as f:
                    f.write(img_resp.content)
                return str(cached_path)
        except Exception:
            return None

    def _render_ppt_from_plan(self, plan: Dict[str, Any]) -> str:
        prs = Presentation()
        out_dir = "outputs"
        os.makedirs(out_dir, exist_ok=True)
        sections = plan.get("content_plan", {}).get("sections", [])

        BG_COLOR = RGBColor(230, 235, 242) 
        TITLE_COLOR = RGBColor(10, 50, 100) 
        TEXT_COLOR = RGBColor(50, 60, 70)  
        ACCENT_COLOR = RGBColor(220, 80, 50) # Orange

        for i, section in enumerate(sections):
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = BG_COLOR

            if i == 0:
                # --- TITLE SLIDE (ORANGE BAR ON TOP) ---
                accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
                accent.fill.solid()
                accent.fill.fore_color.rgb = ACCENT_COLOR
                accent.line.fill.background()

                title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(2.0))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = section.get("title", "Presentation Title")
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = TITLE_COLOR
                p.alignment = 1 # Center

                desc = section.get("description", "").strip()
                if desc:
                    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.5))
                    p_d = desc_box.text_frame.paragraphs[0]
                    p_d.text = desc
                    p_d.font.size = Pt(22)
                    p_d.font.italic = True
                    p_d.font.color.rgb = TEXT_COLOR
                    p_d.alignment = 1
                continue

            # --- CONTENT SLIDES (ORANGE BAR ON LEFT) ---
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.1), Inches(0.6))
            accent.fill.solid()
            accent.fill.fore_color.rgb = ACCENT_COLOR
            accent.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(9.0), Inches(0.8))
            p_t = title_box.text_frame.paragraphs[0]
            p_t.text = section.get("title", "Section Title")
            p_t.font.size = Pt(32)
            p_t.font.bold = True
            p_t.font.color.rgb = TITLE_COLOR 

            img_path = self._fetch_image_for_query(section.get("image_search_query", "modern"))
            text_width = Inches(5.0) if img_path else Inches(9.0)
            if img_path:
                try:
                    slide.shapes.add_picture(img_path, Inches(5.8), Inches(1.5), width=Inches(3.8))
                except Exception: pass

            y_pos = 1.3
            desc = section.get("description", "").strip()
            if desc:
                d_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), text_width, Inches(1.2))
                d_box.text_frame.word_wrap = True
                p_d = d_box.text_frame.paragraphs[0]
                p_d.text = desc
                p_d.font.size = Pt(16)
                p_d.font.color.rgb = TEXT_COLOR 
                y_pos += 1.0

            elements = section.get("elements", [])
            if elements:
                box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), text_width, Inches(5.0))
                tf = box.text_frame
                tf.word_wrap = True
                for idx, el in enumerate(elements):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    val = el if isinstance(el, str) else el.get("description", el.get("text", ""))
                    p.text = f"• {val}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = TEXT_COLOR
                    p.space_after = Pt(8)

        ts = datetime.utcnow().strftime("%Y%m%dT%H%M%SZ")
        path = os.path.join(out_dir, f"presentation_{ts}.pptx")
        prs.save(path)
        return path
        
    def _render_xlsx_from_plan(self, plan: Dict[str, Any]) -> str:
        wb = Workbook()
        out_dir = "outputs"
        os.makedirs(out_dir, exist_ok=True)
        ts = datetime.utcnow().strftime("%Y%m%dT%H%M%SZ")
        path = os.path.join(out_dir, f"spreadsheet_{ts}.xlsx")
        wb.save(path)
        return path

    def get_tools(self) -> dict[str, Any]:
        return {'generate_document': self.generate_document}